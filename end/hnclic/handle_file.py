import sys, os, zipfile,re, requests,shutil,json,glob
sys.path.append(os.path.realpath(os.curdir+"/hello_app/"))
sys.path.append(os.path.realpath(os.curdir+"/hnclic/"))
import numpy as np
import pandas as pd
import pandasql,excel2img
from numpy import nan as NaN
from functools import reduce
import jinja2,datetime,chardet,time,traceback,locale,copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE#MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT
from pptx.chart.data import ChartData
import comtypes.client
from openpyxl.formula.translate import Translator
from openpyxl  import load_workbook
import asyncio
import aiohttp
from utils import unzip_single,zipDir,get_jinja2_Environment,is_number,exec_template,guess_col_names
locale.setlocale(locale.LC_CTYPE, 'chinese')

def convert_file_for_txt(out_filename,template_file,ds_dict):
    '''按模板转换文本文件
    '''
    if not os.path.exists(os.path.split((os.path.realpath(out_filename)))[0]):
        os.makedirs(os.path.split((os.path.realpath(out_filename)))[0])

    with open(template_file, 'rb') as f:
        data = f.read()
        f_charInfo=chardet.detect(data)
    encoding='utf-8' if f_charInfo['encoding'] is None else f_charInfo['encoding']        
    with open(template_file,mode='r',encoding=encoding) as fr \
        ,open(out_filename,mode='w',encoding=encoding) as fw:
        res = fr.read()
        result=exec_template(None,res,ds_dict)
        fw.seek(0)
        fw.truncate()
        fw.write(result)
        return result

def _get_remark_define(notes_text,ds_dict):
    ds_iter,title_lines,Need_lines,loop_var=None,0,20,''
    for one in notes_text.split("\n"):
        pos=one.find('=')
        if pos >-1:#取各种用户定义的变量
            if one[:pos].strip()=='title_lines':
                title_lines=int(one[pos+1:])
            elif one[:pos].strip()=='loop_var':
                loop_var=[x.strip() for x in one[pos+1:].split(',')]
            elif one[:pos].strip()=='Need_lines':
                Need_lines=int(one[pos+1:])
            elif one[:pos].strip()=='dataset':#取dataFrame,切片，取游标
                if len(loop_var)==2:
                    ds_iter=pd.eval(one[pos+1:].strip(),local_dict=ds_dict)[:Need_lines].reset_index(drop=True).iterrows()
                else:
                    ds_iter=pd.eval(one[pos+1:].strip(),local_dict=ds_dict)[:Need_lines].reset_index(drop=True)
    return ds_iter,title_lines,Need_lines,loop_var


def convert_file_for_pptx(out_filename,template_file,ds_dict):
    '''按模板转换xlsx文件
    按字典转换模板文件，输出为out_filename
    '''
    unzip_path=os.path.join(out_filename + 't\\pptx_tmp')
    if(os.path.exists(unzip_path)):
        shutil.rmtree(unzip_path)
    unzip_single(template_file,unzip_path)
    embeddings_path=os.path.join(unzip_path,"ppt\\embeddings")
    tmp_pd_dict={}
    if(os.path.exists(embeddings_path)):
        for x in os.listdir(embeddings_path):
            if x.endswith('.xlsx'):
                convert_file_for_xlsx(os.path.join(embeddings_path,x),os.path.join(embeddings_path,x),ds_dict)
                tmp_pd_dict[x]=pd.read_excel(os.path.join(embeddings_path,x))
    zipDir(unzip_path,out_filename)
    shutil.rmtree(out_filename+"t")

    env = get_jinja2_Environment()
    ppt_file = Presentation(out_filename)

    '''
    #expr
    title_lines=1
    loop_var=index,row
    dataset=a.sort_values(zhibiao,ascending=False)[:size]
    '''
    def calc_frame_txt(obj,calc_dict,calc_kind=None):
        if calc_kind is None:
            calc_kind=1 if len(obj.text_frame.paragraphs)<1  else 3
        if calc_kind==3:#text_frame 中有多个不同格式的文本，需要查runs,通常不应该是这样的
            for paragraph in obj.text_frame.paragraphs:
                exp_list=[]
                if paragraph.text.find('{{')>-1:
                    start,end,s_num,e_num=-1,-1,0,0
                    for idx,run in enumerate(paragraph.runs):
                        if run.text.find('{{')>-1:
                            s_num+=1
                            if s_num==1:
                                start=idx
                        if run.text.find('}}')>-1:
                            end=idx
                            e_num+=1
                        if start>=0 and end >=0 and s_num==e_num:
                            exp_list.append((start,end))
                            start,end,s_num,e_num=-1,-1,0,0
                    for start,end in exp_list:
                        if start>=0 and end >=0 and start<=end:
                            text=''.join([x.text for x in paragraph.runs[start:end+1]])
                            try:
                                result = exec_template(env,text,calc_dict)                                
                            except Exception as e:
                                raise RuntimeError(text)
                            paragraph.runs[start].text=result
                            for x in paragraph.runs[start+1:end+1]:x.text=''
        elif calc_kind==2:
            for paragraph in obj.text_frame.paragraphs:
                if paragraph.text.find('{{')>-1:
                    try:
                        result = exec_template(env,paragraph.text,calc_dict)  
                    except:
                        raise RuntimeError(paragraph.text)                        
                    for run in paragraph.runs:
                        run.text=''#直接copy font 报错，我们通过将其他runs中的文字清空，计算出的新文字赋值给第一个run。这样就保留了格式
                    paragraph.runs[0].text=result
        else:
            expr=obj.text_frame.text
            if expr.find('{{')>-1:
                try:
                    result = exec_template(env,expr,calc_dict) # env.from_string(expr)
                except:
                    raise RuntimeError(paragraph.text)
                for paragraph in obj.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text=''#直接copy font 报错，我们通过将其他runs中的文字清空，计算出的新文字赋值给第一个run。这样就保留了格式
                obj.text_frame.paragraphs[0].runs[0].text=result            
    def handle_all_shapes(shapes,real_dict,tmp_pd_dict):
        # real_dict 我们使用这个参数来层层传递外面定义变量
        #tmp_pd_dict 是专为内嵌excel准备的，貌似递归取不到外层定义的变量
        for shape in shapes:#shape.part.related_parts['rId4'].blob
            if hasattr(shape,"shapes"):
                handle_all_shapes(shape.shapes,real_dict,tmp_pd_dict)
                continue
            if shape.has_text_frame or shape.has_table:
                pass
            if shape.shape_type==MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT :
                pass
            if shape.has_text_frame:
                calc_frame_txt(shape,real_dict)
            elif shape.has_chart:
                key=shape.chart._workbook.xlsx_part.partname.split("/")[-1]
                # 定义图表数据 ---------------------
                chart_data = ChartData()
                columns=list(tmp_pd_dict[key].columns.values)
                chart_data.categories =tmp_pd_dict[key][columns[0]]
                for one in columns[1:]:
                    chart_data.add_series(one, tuple(tmp_pd_dict[key][one]))
                shape.chart.replace_data(chart_data)
            elif shape.has_table:
                current_row=0
                for row in shape.table.rows:
                    current_col=0
                    for cell in row.cells:
                        if cell.text_frame.text.find('{{')<0:
                            current_col=current_col+1
                            continue
                        try:
                            result=exec_template(env,cell.text_frame.text,real_dict)
                        except:
                            raise RuntimeError(cell.text_frame.text)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.text=''#直接copy font 报错，我们通过将其他runs中的文字清空，计算出的新文字赋值给第一个run。这样就保留了格式                        
                        copy_row=current_row                                
                        result_lines=result.split('\n')
                        for one_line in result_lines:#展开模板计算结果
                            copy_col=current_col
                            #从当前位置开始，复制结果到ppt的table中
                            for one in one_line.split():
                                cur_row_cells=shape.table.rows[copy_row].cells
                                if copy_col>=len(cur_row_cells):#如果ppt table 中的列不够用，当前行的复制就结束
                                    break
                                p_cell=cur_row_cells[copy_col]
                                if len(p_cell.text_frame.paragraphs[0].runs)==0:
                                    p_cell.text_frame.paragraphs[0].add_run()
                                p_cell.text_frame.paragraphs[0].runs[0].text=one
                                copy_col=copy_col+1
                            copy_row=copy_row+1
                            if copy_row>=len(shape.table.rows):#行不够就结束复制
                                break
                        current_col=current_col+1
                    current_row=current_row+1
                    if current_row>=len(shape.table.rows):
                        break    
    try:
        real_dict=ds_dict.copy()
        for slide in ppt_file.slides:
            if slide.has_notes_slide:#抽取备注栏里面的变量定义，后页会覆盖前页
                notes_text=slide.notes_slide.notes_text_frame.text
                for one_line in notes_text.split("\n"):
                    var_expr=one_line.split("=")
                    if len(var_expr)<2:
                        continue
                    try:
                        if var_expr[1].strip().startswith("{{"):
                            result_lines = exec_template(env,var_expr[1],real_dict)
                        else:
                            result_lines = exec_template(env,"{{" + var_expr[1]+ "}}",real_dict)
                        real_dict=real_dict.copy()
                        real_dict[var_expr[0].strip() ] = result_lines
                    except Exception as e:
                        raise RuntimeError("\n备注说明中的公式不正确："+one_line)       

            handle_all_shapes(slide.shapes,real_dict,tmp_pd_dict) 
    
        ppt_file.save(out_filename)
    finally:
        if ppt_file is not None:
            ppt_file.save(out_filename)
            del ppt_file
        ppt2png(out_filename,ds_dict.get("_idx_",''))

def ppt2png(pptFileName,idx=0):
    comtypes.CoInitialize()
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = True
        outputFileName = pptFileName[0:-5] + ".pdf"
        ppt = powerpoint.Presentations.Open(pptFileName)
        try:
            #保存为图片
            ppt.SaveAs(pptFileName[0:-5] + f'_{idx}.jpg', 17)
            #保存为pdf
            #ppt.SaveAs(outputFileName, 32) # formatType = 32 for ppt to pdf
            # 关闭打开的ppt文件
        finally:
            # 关闭打开的ppt文件
            if ppt is not None:
                ppt.Close()
                ppt=None
            # 关闭powerpoint软件
            if powerpoint is not None:
                powerpoint.Quit()
                powerpoint=None
    finally:
        comtypes.CoUninitialize()



def convert_file_for_xlsx(out_filename,template_file,ds_dict, appendFunDict=None):
    '''按模板转换xlsx文件
    按字典转换模板文件，输出为out_filename
    '''
    def calc_cell(sheet):
        for row in sheet.rows:
            need_lines=0
            for cell in row:
                #模板计算
                if cell.value is not None and cell.data_type=='s' and cell.value.find('{{')>-1 :
                    result = exec_template(env,cell.value,real_dict)  
                    row=cell.row
                    start_col=cell.column
                    result_lines=result.split('\n')
                    for one_line in result_lines:#展开模板计算结果
                        col=start_col
                        for one in one_line.split():
                            p_cell=sheet.cell(row=row,column=col)
                            if len(one)<14 and is_number(one) :#14位工号，已经到万亿了，现在还不可能有这么大的数
                                p_cell.value=float(one) 
                                p_cell.data_type='n'
                            else:
                                p_cell.value=one
                            col=col+1
                        row=row+1
                    need_lines=row-cell.row
                    continue
                #复制公式
                elif cell.value is not None and cell.data_type=='f' and cell.value.startswith('=') :
                    row=cell.row+1
                    for one in range(1,need_lines):
                        p_cell=sheet.cell(row=row,column=cell.column).coordinate
                        sheet[p_cell] = Translator(cell.value, origin=cell.coordinate).translate_formula(p_cell)
                        row=row+1
                    need_lines= row-cell.row
                    continue
                need_lines= need_lines

    env=get_jinja2_Environment()
    wb=None
    try:
        wb = load_workbook(template_file)
        real_dict=ds_dict.copy()
        
        #loop;单位;a
        list_g=[]
        for sheet in wb.worksheets:
            if sheet.title.startswith('loop'):
                groupby=sheet.title.split(";")[1]
                ds_name=sheet.title.split(";")[2]
                list_g=appendFunDict['list_group'](groupby)
                sheet.title=list_g[0]
                for i in range(1,len(list_g)):
                    copy_sheet=wb.copy_worksheet(sheet)
                    copy_sheet.title=list_g[i]

                for one in list_g:
                    real_dict=ds_dict.copy()
                    real_dict[ds_name]=appendFunDict['groupby'](one)
                    calc_cell(wb[one])

        real_dict=ds_dict.copy()
        for sheet in wb.worksheets:
                calc_cell(sheet)
        wb.save(out_filename)
        cnt=0
        page_images=[]
        for sheet in wb.worksheets: 
            if sheet.title.startswith("_") or (list_g is not None and list_g.count(sheet.title)>0):
                page_images.append((f"{out_filename}{ '{:0>2d}'.format(cnt) }{sheet.title}.png", sheet.title))
            cnt=cnt+1
        excel2img.export_img_many(out_filename,page_images)
    finally:
        if wb!=None:
            wb.close()
            del wb
            #excel_catch_screen(out_filename,sheet.title,sheet.dimensions,f"{out_filename}{sheet.title}.png")

def json_keyvalue_all(input_json,rootStack=[],ds_dict={}):
    try:
        env = get_jinja2_Environment()
        ret_val=False
        if isinstance(input_json,dict):
            for key in input_json.keys():
                key_value = input_json.get(key)
                if isinstance(key_value,dict):
                    rootStack.append((input_json,key ))
                    return json_keyvalue_all(key_value,rootStack,ds_dict)
                elif isinstance(key_value,list):
                    for idx,json_array in enumerate(key_value):
                        rootStack.append((key_value,idx ))
                        c_ret=ret_val or json_keyvalue_all(json_array,rootStack,ds_dict)
                        ret_val=ret_val or c_ret
                        if ret_val:
                            #ar2 = list(map(list,zip(*arr))) #行列互换
                            arr=[one.split('\n') for one in json_array]
                            max_len=max([ len(x) for x in arr])
                            input_json[key]= key_value[0:idx] + [ [row[i] if len(row)>i else '' for row in arr] for i in range(max_len) ]
                            break                        
                else:
                    if isinstance(key_value,str) and key_value.find("{{")>=0:#转换
                        print (str(key)+ " = " +str(key_value))
                        input_json[key]=exec_template(env,key_value,ds_dict)#"被替换了a"
                        #return True
        elif isinstance(input_json,list):
            for idx,input_json_array in enumerate(input_json):
                rootStack.append((input_json,idx ))
                c_ret=json_keyvalue_all(input_json_array,rootStack,ds_dict)
                ret_val=ret_val or c_ret
            return ret_val
        else:
            if isinstance(input_json,str) and input_json.find("{{")>=0:#转换
                print ( " 数组内： " +str(input_json))
                val=exec_template(env,input_json,ds_dict) #"被替\n换了c" 
                if isinstance(rootStack[-1][0],list) :
                    rootStack[-1][0][ rootStack[-1][1] ]=val
                elif isinstance(rootStack[-1],dict) :
                    rootStack[-1][0][ rootStack[-1][1] ]=val
                return True
    finally:
        rootStack.pop()
    return False

def convert_html(fileName,url,ds_dict):
    html_text=requests.get(url).text
    start_pos=html_text.find("<script>window.__work =")+len("<script>window.__work =")
    end_pos=html_text.find("</script>",start_pos)
    _json_str=html_text[start_pos:end_pos]
    work_json=json.loads(_json_str)
    json_keyvalue_all(work_json,[(work_json,'')],ds_dict)
    html_text=html_text[0:start_pos]+ json.dumps(work_json) +html_text[end_pos:]
    html_text=html_text.replace("<head>",f'<head><base href="{url}" />')
    fo = open(fileName, "w",encoding='utf8')
    fo.write(html_text)
    fo.close()
 